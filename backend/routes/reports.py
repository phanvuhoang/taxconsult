import asyncio
import io
import json
import re
import time
import uuid

from fastapi import APIRouter, Depends, HTTPException, Query, Body
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select, desc
from typing import Optional, List

from backend.database import get_db, AsyncSessionLocal, DbvntaxSession
from backend.models import Report, ReportJob, User
from backend.auth import get_current_user
from backend.ai_provider import call_ai
from backend.config import DEFAULT_SECTIONS as CONFIG_SECTIONS, SECTOR_SECTIONS, COMPANY_SECTIONS

router = APIRouter(prefix="/api/reports", tags=["reports"])


class FullReportRequest(BaseModel):
    subject: str
    mode: str = "ngành"
    tax_types: List[str] = ["TNDN", "GTGT"]
    time_period: Optional[str] = None
    model_tier: str = "deepseek"
    sonar_model: str = "sonar"
    sections: Optional[List[dict]] = None


class ExportRequest(BaseModel):
    subject: str
    html_content: str


class GammaRequest(BaseModel):
    subject: str
    html_content: str
    num_cards: int = 20


# ── List / Get / Delete (unchanged) ──────────────────────────────────────────

@router.get("")
async def list_reports(
    report_type: str = Query(None),
    skip: int = 0,
    limit: int = 20,
    user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    q = select(Report).order_by(desc(Report.created_at))
    if user.role != "admin":
        q = q.where(Report.user_id == user.id)
    if report_type:
        q = q.where(Report.report_type == report_type)
    q = q.offset(skip).limit(limit)
    result = await db.execute(q)
    reports = result.scalars().all()
    return [
        {
            "id": r.id,
            "title": r.title,
            "subject": r.subject,
            "report_type": r.report_type,
            "tax_types": r.tax_types,
            "time_period": r.time_period,
            "model_used": r.model_used,
            "duration_ms": r.duration_ms,
            "created_at": r.created_at.isoformat() if r.created_at else None,
        }
        for r in reports
    ]


@router.get("/job/{job_id}")
async def get_job_status(
    job_id: str,
    user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    job = await db.get(ReportJob, job_id)
    if not job:
        raise HTTPException(status_code=404, detail="Job not found")
    citations = []
    if job.status == "done" and job.report_id:
        saved = await db.get(Report, job.report_id)
        if saved and saved.citations:
            citations = saved.citations

    return {
        "status": job.status,
        "progress_step": job.progress_step,
        "progress_total": job.progress_total,
        "progress_label": job.progress_label,
        "html_content": job.html_content,
        "error_msg": job.error_msg,
        "report_id": job.report_id,
        "citations": citations,
    }


@router.post("/start")
async def start_report(
    body: FullReportRequest,
    user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Tạo job, spawn background task, trả về job_id ngay."""
    if not body.subject.strip():
        raise HTTPException(status_code=400, detail="Subject is required")

    job_id = str(uuid.uuid4())
    sections = body.sections if body.sections else CONFIG_SECTIONS
    enabled_count = len([s for s in sections if s.get("enabled", True)])

    job = ReportJob(
        id=job_id,
        subject=body.subject,
        user_id=user.id,
        status="pending",
        progress_total=enabled_count * 2,
    )
    db.add(job)
    await db.commit()

    # Normalize sections in body before passing to background task
    body_with_sections = FullReportRequest(
        subject=body.subject,
        mode=body.mode,
        tax_types=body.tax_types,
        time_period=body.time_period,
        model_tier=body.model_tier,
        sonar_model=body.sonar_model,
        sections=sections,
    )

    # Fire and forget — không await
    asyncio.create_task(run_report_job(job_id, body_with_sections, user.id))
    return {"job_id": job_id}


@router.get("/{report_id}")
async def get_report(
    report_id: int,
    user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    result = await db.execute(select(Report).where(Report.id == report_id))
    report = result.scalar_one_or_none()
    if not report:
        raise HTTPException(status_code=404, detail="Report not found")
    if user.role != "admin" and report.user_id != user.id:
        raise HTTPException(status_code=403, detail="Forbidden")
    return {
        "id": report.id,
        "title": report.title,
        "subject": report.subject,
        "report_type": report.report_type,
        "tax_types": report.tax_types,
        "time_period": report.time_period,
        "content_html": report.content_html,
        "citations": report.citations,
        "model_used": report.model_used,
        "duration_ms": report.duration_ms,
        "created_at": report.created_at.isoformat() if report.created_at else None,
    }


@router.delete("/{report_id}")
async def delete_report(
    report_id: int,
    user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    result = await db.execute(select(Report).where(Report.id == report_id))
    report = result.scalar_one_or_none()
    if not report:
        raise HTTPException(status_code=404, detail="Report not found")
    if user.role != "admin" and report.user_id != user.id:
        raise HTTPException(status_code=403, detail="Forbidden")
    await db.delete(report)
    await db.commit()
    return {"ok": True}


@router.get("/{report_id}/export-docx")
async def export_docx_by_id(
    report_id: int,
    user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    result = await db.execute(select(Report).where(Report.id == report_id))
    report = result.scalar_one_or_none()
    if not report:
        raise HTTPException(status_code=404, detail="Report not found")
    if user.role != "admin" and report.user_id != user.id:
        raise HTTPException(status_code=403, detail="Forbidden")

    docx_bytes = _html_to_docx(report.content_html or "", report.title)
    filename = f"report_{report_id}.docx"
    return StreamingResponse(
        io.BytesIO(docx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@router.post("/docx")
async def export_docx(
    body: ExportRequest,
    user: User = Depends(get_current_user),
):
    """Export HTML content to DOCX file."""
    docx_bytes = _html_to_docx(body.html_content, body.subject)
    safe_name = re.sub(r'[^\w\s-]', '', body.subject)[:60].strip().replace(' ', '_')
    filename = f"{safe_name}.docx"
    return StreamingResponse(
        io.BytesIO(docx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@router.post("/slides")
async def export_slides(
    body: ExportRequest,
    user: User = Depends(get_current_user),
):
    """Export HTML content to PPTX slides."""
    pptx_bytes = _html_to_slides(body.html_content, body.subject)
    safe_name = re.sub(r'[^\w\s-]', '', body.subject)[:60].strip().replace(' ', '_')
    filename = f"{safe_name}.pptx"
    return StreamingResponse(
        io.BytesIO(pptx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


@router.post("/gamma")
async def create_gamma(
    body: GammaRequest,
    user: User = Depends(get_current_user),
):
    """Tạo Gamma presentation từ report HTML."""
    import os
    import httpx as _httpx
    gamma_key = os.getenv("GAMMA_API_KEY", "")
    if not gamma_key:
        raise HTTPException(status_code=400, detail="GAMMA_API_KEY not configured")

    from backend.doc_context import strip_html_tvpl
    text_content = strip_html_tvpl(body.html_content)

    async with _httpx.AsyncClient(timeout=60) as client:
        r = await client.post(
            "https://api.gamma.app/v1/presentations",
            headers={
                "Authorization": f"Bearer {gamma_key}",
                "Content-Type": "application/json",
            },
            json={
                "title": f"Phân tích thuế: {body.subject}",
                "text": text_content[:8000],
                "num_cards": min(body.num_cards, 60),
                "theme": "professional",
                "language": "vi",
            },
        )
        try:
            r.raise_for_status()
        except Exception:
            raise HTTPException(status_code=502, detail=f"Gamma API error: {r.text[:200]}")
        data = r.json()
        return {"url": data.get("url"), "id": data.get("id")}


@router.post("/suggest-subsections")
async def suggest_subsections(
    body: dict = Body(...),
    user: User = Depends(get_current_user),
):
    """Gợi ý sub-topics cho một section cụ thể."""
    title = body.get("title", "")
    subject = body.get("subject", "")
    prompt = (
        f'Đề xuất 4-5 chủ đề con cho phần "{title}" '
        f'trong báo cáo thuế về: {subject}.\n'
        f'Trả về JSON array. Chỉ trả về array, không giải thích.\n'
        f'Ví dụ: ["Chủ đề 1", "Chủ đề 2"]'
    )
    result = await call_ai(
        messages=[{"role": "user", "content": prompt}],
        model_tier="haiku",
        max_tokens=400,
    )
    match = re.search(r'\[.*?\]', result["content"], re.DOTALL)
    suggestions = json.loads(match.group()) if match else []
    return {"suggestions": suggestions}


@router.get("/default-sections")
async def get_default_sections(
    mode: str = Query("ngành"),
    user: User = Depends(get_current_user),
):
    return COMPANY_SECTIONS if mode == "công ty" else SECTOR_SECTIONS


@router.post("/suggest-topics")
async def suggest_topics(
    subject: str = Body(...),
    mode: str = Body("ngành"),
    tax_types: list = Body([]),
    current_user: User = Depends(get_current_user),
):
    """Dùng Claudible Haiku để gợi ý sections và sub-topics."""
    system = "Bạn là chuyên gia tư vấn thuế Việt Nam. Trả lời bằng JSON."
    prompt = f"""Tôi cần viết báo cáo phân tích thuế cho: "{subject}" (loại: {mode}).
Các sắc thuế quan tâm: {', '.join(tax_types) or 'tổng quát'}.

Hãy gợi ý danh sách sections và sub-topics phù hợp nhất cho báo cáo này.
Trả về JSON với format:
{{
  "sections": [
    {{"id": "s1", "title": "Tên section", "enabled": true, "sub": ["sub-topic 1", "sub-topic 2"]}},
    ...
  ]
}}
Tối đa 8 sections, mỗi section tối đa 5 sub-topics. Ưu tiên các vấn đề thuế đặc thù của ngành/chủ đề."""

    result = await call_ai(
        messages=[{"role": "user", "content": prompt}],
        system=system,
        model_tier="haiku",
        max_tokens=2000,
    )
    content = result["content"]
    json_match = re.search(r'\{.*\}', content, re.DOTALL)
    if json_match:
        data = json.loads(json_match.group())
        return data
    return {"sections": []}


# ── Background worker ─────────────────────────────────────────────────────────

async def run_report_job(job_id: str, body: FullReportRequest, user_id: int):
    """Chạy hoàn toàn background — không phụ thuộc client connection."""
    from backend.report_generator import (
        _gather_section_context,
        SECTION_PROMPT_TAX,
        SECTION_PROMPT_GENERAL,
        SECTION_SYSTEM,
    )
    from backend.time_period import parse_time_period
    from backend.doc_context import get_priority_docs_context, get_priority_doc_ids

    async with AsyncSessionLocal() as db:
        async with DbvntaxSession() as dbvntax_db:
            job = await db.get(ReportJob, job_id)
            job.status = "running"
            await db.commit()

            try:
                sections_config = body.sections if body.sections else CONFIG_SECTIONS
                enabled = [s for s in sections_config if s.get("enabled", True)]
                period = parse_time_period(body.time_period or "")

                job.progress_total = len(enabled) * 2
                job.progress_label = "Đang thu thập dữ liệu..."
                await db.commit()

                # Phase 0: priority docs context (shared across sections)
                priority_ctx = ""
                exclude_ids = []
                try:
                    priority_ctx = await get_priority_docs_context(
                        db, dbvntax_db, body.tax_types,
                        time_period_end=period["end_date"],
                        time_period_start=period.get("start_date"),
                    )
                    exclude_ids = await get_priority_doc_ids(db, body.tax_types)
                except Exception:
                    pass

                # Phase 1: gather per-section contexts in parallel (perplexity + docs)
                job.progress_step = 0
                job.progress_label = "Đang research tất cả sections..."
                await db.commit()

                context_tasks = [
                    _gather_section_context(
                        sec, body.subject, body.tax_types, period,
                        dbvntax_db, body.sonar_model,
                        exclude_dbvntax_ids=exclude_ids if exclude_ids else None,
                    )
                    for sec in enabled
                ]
                all_contexts = await asyncio.gather(*context_tasks, return_exceptions=True)

                # Phase 2: write each section with AI
                full_html = f"<h1>Báo cáo Phân tích Thuế: {body.subject}</h1>\n"
                full_html += f"<p><em>Giai đoạn: {period['label']} | Sắc thuế: {', '.join(body.tax_types)}</em></p>\n"

                start_time = time.time()
                all_citations: list = []
                for ctx in all_contexts:
                    if isinstance(ctx, dict):
                        all_citations.extend(ctx.get("citations", []))
                # deduplicate preserving order
                seen = set()
                deduped_citations = []
                for c in all_citations:
                    if c not in seen:
                        seen.add(c)
                        deduped_citations.append(c)

                for idx, (section, ctx) in enumerate(zip(enabled, all_contexts)):
                    if isinstance(ctx, Exception):
                        ctx = {"perplexity": "", "docs": "", "congvan": ""}

                    job.progress_step = len(enabled) + idx
                    job.progress_label = f"AI đang viết: {section.get('title', '')}"
                    await db.commit()

                    section_number = idx + 1
                    if section.get("tax_aware"):
                        prompt = SECTION_PROMPT_TAX.format(
                            section_title=section["title"],
                            section_number=section_number,
                            subject=body.subject,
                            mode=body.mode,
                            time_period=period["label"],
                            priority_context=priority_ctx or "(Không có văn bản ưu tiên)",
                            tax_docs_context=ctx.get("docs", "") or "(Không có dữ liệu)",
                            congvan_context=ctx.get("congvan", "") or "(Không có dữ liệu)",
                            perplexity_context=ctx.get("perplexity", "") or "(Không có dữ liệu)",
                        )
                    else:
                        prompt = SECTION_PROMPT_GENERAL.format(
                            section_title=section["title"],
                            section_number=section_number,
                            subject=body.subject,
                            mode=body.mode,
                            time_period=period["label"],
                            perplexity_context=ctx.get("perplexity", "") or "(Không có dữ liệu)",
                        )

                    result = await call_ai(
                        messages=[{"role": "user", "content": prompt}],
                        system=SECTION_SYSTEM,
                        model_tier=body.model_tier,
                        max_tokens=8192,
                    )
                    section_html = result["content"] if result else f"<h2>{section_number}. {section['title']}</h2><p>(Không thể tạo nội dung)</p>"
                    full_html += section_html + "\n"

                    # Update partial HTML so frontend can see progress
                    job.html_content = full_html
                    await db.commit()

                # Save completed report to DB
                duration_ms = int((time.time() - start_time) * 1000)
                report = Report(
                    user_id=user_id,
                    title=f"Phân tích thuế: {body.subject}",
                    subject=body.subject,
                    report_type="full",
                    tax_types=body.tax_types,
                    time_period=body.time_period,
                    content_html=full_html,
                    citations=deduped_citations,
                    model_used=body.model_tier,
                    duration_ms=duration_ms,
                )
                db.add(report)
                await db.commit()
                await db.refresh(report)

                job.status = "done"
                job.html_content = full_html
                job.progress_label = "Hoàn thành!"
                job.report_id = report.id
                await db.commit()

            except Exception as e:
                job.status = "error"
                job.error_msg = str(e)
                await db.commit()


# ── HTML → DOCX ───────────────────────────────────────────────────────────────

def _html_to_docx(html: str, title: str) -> bytes:
    """Convert HTML to DOCX using python-docx + BeautifulSoup."""
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.oxml.ns import qn
    from bs4 import BeautifulSoup

    BRAND_COLOR = RGBColor(0x02, 0x8a, 0x39)

    doc = Document()

    # Style the title
    title_para = doc.add_heading(title, 0)
    for run in title_para.runs:
        run.font.color.rgb = BRAND_COLOR

    soup = BeautifulSoup(html, "html.parser")

    for elem in soup.find_all(["h1", "h2", "h3", "h4", "p", "li", "table"]):
        tag = elem.name
        text = elem.get_text(separator=" ", strip=True)
        if not text:
            continue
        if tag == "h1":
            h = doc.add_heading(text, level=1)
            for run in h.runs:
                run.font.color.rgb = BRAND_COLOR
        elif tag == "h2":
            h = doc.add_heading(text, level=2)
            for run in h.runs:
                run.font.color.rgb = BRAND_COLOR
        elif tag == "h3":
            doc.add_heading(text, level=3)
        elif tag == "h4":
            doc.add_heading(text, level=4)
        elif tag == "li":
            doc.add_paragraph(text, style="List Bullet")
        elif tag == "table":
            rows = elem.find_all("tr")
            if not rows:
                continue
            cols = max(len(r.find_all(["td", "th"])) for r in rows)
            if cols == 0:
                continue
            table = doc.add_table(rows=len(rows), cols=cols)
            table.style = "Table Grid"
            for ri, row in enumerate(rows):
                cells = row.find_all(["td", "th"])
                for ci, cell in enumerate(cells):
                    if ci < cols:
                        table.cell(ri, ci).text = cell.get_text(strip=True)
        else:
            doc.add_paragraph(text)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── HTML → PPTX ───────────────────────────────────────────────────────────────

def _html_to_slides(html: str, title: str) -> bytes:
    """Convert HTML to PPTX slides using python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from bs4 import BeautifulSoup

    BRAND = RGBColor(0x02, 0x8a, 0x39)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Title slide
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    for para in slide.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = BRAND
            run.font.bold = True
            run.font.size = Pt(32)
    try:
        slide.placeholders[1].text = "Báo cáo Phân tích Thuế"
    except Exception:
        pass

    soup = BeautifulSoup(html, "html.parser")

    current_title = None
    current_bullets: list = []

    def flush_slide():
        if not current_title:
            return
        layout = prs.slide_layouts[1]  # Title and Content
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = current_title
        for para in s.shapes.title.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = BRAND
                run.font.bold = True
                run.font.size = Pt(24)
        if current_bullets:
            try:
                tf = s.placeholders[1].text_frame
                tf.clear()
                for i, bullet in enumerate(current_bullets[:12]):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = bullet[:200]
                    p.font.size = Pt(14)
            except Exception:
                pass

    for elem in soup.find_all(["h1", "h2", "h3", "p", "li"]):
        tag = elem.name
        text = elem.get_text(separator=" ", strip=True)
        if not text:
            continue
        if tag in ("h1", "h2"):
            flush_slide()
            current_title = text
            current_bullets = []
        elif current_title and text:
            current_bullets.append(text)

    flush_slide()

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
